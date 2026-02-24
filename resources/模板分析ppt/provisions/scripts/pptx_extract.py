#!/usr/bin/env python3
"""
PPTX 提取模块

从 PPTX 文件中提取幻灯片数据为字典格式。
支持所有形状类型：文本、图片、表格、图表、组合形状、自动形状等。

关键设计：
1. 为所有形状保存原始 XML（包括 TEXT_BOX），确保重建时保留完整格式
2. 为含图片引用的形状保存关系映射（xml_rels），确保重建时能恢复图片
3. 图片数据去重存储在顶层 images 字典中
"""

import base64
import hashlib
import subprocess
import sys
import logging
from pathlib import Path
from typing import Dict, Any, Optional, List
from copy import deepcopy

logger = logging.getLogger(__name__)


def ensure_dependencies() -> bool:
    """确保依赖已安装"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from lxml import etree
        return True
    except ImportError:
        try:
            subprocess.check_call([
                sys.executable, "-m", "pip", "install",
                "python-pptx>=0.6.21", "lxml", "--quiet"
            ])
            return True
        except Exception as e:
            logger.error(f"无法安装依赖: {e}")
            return False


R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
_RID_ATTRS = [f'{{{R_NS}}}embed', f'{{{R_NS}}}link', f'{{{R_NS}}}id']


def _extract_shape_xml(shape) -> Optional[str]:
    """
    提取形状的原始 XML。

    对所有形状类型都保存——这确保重建时能完整保留格式
    （填充色、边框、阴影、字体样式等），不仅仅是内容。
    """
    try:
        from lxml import etree
        if hasattr(shape, '_element'):
            return etree.tostring(shape._element, encoding='unicode')
    except Exception as e:
        logger.warning(f"无法提取形状 XML: {e}")
    return None


def _extract_xml_relationships(shape, slide_part, all_images: dict) -> Optional[Dict]:
    """
    提取形状 XML 中引用的关系（图片等）。

    扫描形状 XML 中的所有 r:embed、r:link、r:id 引用，
    找到对应的图片数据，建立 rId → image_ref 映射。

    这使得重建时可以恢复图片关系，解决"组合形状中图片丢失"的问题。

    Args:
        shape: python-pptx Shape 对象
        slide_part: 幻灯片的 Part（包含 rels）
        all_images: 全局图片字典（用于去重存储）

    Returns:
        dict: {rId: {"image_ref": hash}} 或 None
    """
    try:
        from lxml import etree

        if not hasattr(shape, '_element'):
            return None

        # 收集 XML 中所有 rId 引用
        rId_refs = set()
        for elem in shape._element.iter():
            for attr in _RID_ATTRS:
                val = elem.get(attr)
                if val:
                    rId_refs.add(val)

        if not rId_refs:
            return None

        xml_rels = {}
        for rId in rId_refs:
            try:
                rel = slide_part.rels[rId]
                if rel.is_external:
                    continue  # 超链接等外部引用不需要映射

                target_part = rel.target_part
                if not hasattr(target_part, 'blob'):
                    continue

                blob = target_part.blob
                content_type = target_part.content_type

                # 只处理图片类型
                if not content_type.startswith('image/'):
                    continue

                img_hash = hashlib.md5(blob).hexdigest()[:8]

                # 存入全局图片字典（去重）
                if img_hash not in all_images:
                    all_images[img_hash] = {
                        "type": content_type,
                        "data": base64.b64encode(blob).decode('utf-8')
                    }

                xml_rels[rId] = {"image_ref": img_hash}

            except (KeyError, AttributeError):
                continue

        return xml_rels if xml_rels else None

    except Exception as e:
        logger.warning(f"提取关系映射失败: {e}")
        return None


def _try_extract_image(shape, shape_info: Dict) -> bool:
    """
    尝试从形状中提取图片数据。
    检查 PICTURE 类型和占位符中的图片。

    Returns:
        True if image was extracted
    """
    type_name = ""
    try:
        type_name = shape.shape_type.name if hasattr(shape, 'shape_type') else ""
    except Exception:
        pass

    has_image = type_name == "PICTURE"

    if not has_image:
        try:
            if hasattr(shape, 'image') and shape.image is not None:
                _ = shape.image.blob
                has_image = True
        except Exception:
            pass

    if has_image:
        try:
            image = shape.image
            blob = image.blob
            shape_info["image_type"] = image.content_type
            shape_info["image_data"] = base64.b64encode(blob).decode('utf-8')
            shape_info["image_hash"] = hashlib.md5(blob).hexdigest()[:8]
            return True
        except Exception as e:
            shape_info["image_error"] = str(e)
            logger.warning(f"提取图片失败 (name={shape_info.get('name', '?')}): {e}")

    return False


def _extract_text_content(shape, shape_info: Dict) -> bool:
    """提取形状中的文本内容（段落、runs、格式）"""
    if not (hasattr(shape, 'text_frame') and shape.has_text_frame):
        return False

    try:
        paragraphs = []
        for para in shape.text_frame.paragraphs:
            para_info = {
                "text": para.text,
                "level": para.level,
                "runs": []
            }
            try:
                if para.alignment is not None:
                    para_info["alignment"] = str(para.alignment)
            except Exception:
                pass

            for run in para.runs:
                run_info = {"text": run.text}
                try:
                    if run.font.bold:
                        run_info["bold"] = True
                    if run.font.italic:
                        run_info["italic"] = True
                    if run.font.size:
                        run_info["font_size"] = run.font.size
                    if run.font.name:
                        run_info["font_name"] = run.font.name
                    if run.font.color and run.font.color.type is not None:
                        try:
                            if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                run_info["color"] = str(run.font.color.rgb)
                        except (AttributeError, TypeError):
                            pass
                except Exception:
                    pass
                para_info["runs"].append(run_info)
            paragraphs.append(para_info)

        shape_info["text_content"] = paragraphs
        shape_info["full_text"] = shape.text_frame.text
        return bool(shape.text_frame.text.strip())
    except Exception as e:
        logger.warning(f"提取文本失败 (name={shape_info.get('name', '?')}): {e}")
        return False


def _extract_table_data(shape, shape_info: Dict) -> bool:
    """提取表格数据"""
    try:
        if not hasattr(shape, 'table'):
            return False
        table = shape.table
        shape_info["table_data"] = {
            "rows": len(table.rows),
            "columns": len(table.columns),
            "cells": []
        }
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                shape_info["table_data"]["cells"].append({
                    "row": row_idx,
                    "col": col_idx,
                    "text": cell.text
                })
        return True
    except Exception as e:
        logger.warning(f"提取表格失败: {e}")
        return False


def _extract_group_shapes(shape, shape_info: Dict) -> bool:
    """提取组合形状中的子形状"""
    try:
        if not hasattr(shape, 'shapes'):
            return False
        shape_info["group_shapes"] = []
        for i, child in enumerate(shape.shapes):
            child_info = extract_shape_info(child)
            child_info["group_child_index"] = i
            shape_info["group_shapes"].append(child_info)
        return True
    except Exception as e:
        logger.warning(f"提取组合形状失败: {e}")
        return False


def extract_shape_info(shape, slide_part=None, all_images: dict = None) -> Dict[str, Any]:
    """
    提取单个形状的完整信息。

    对所有形状类型：
    1. 提取基础属性（位置、尺寸、名称）
    2. 提取特定内容（文本、图片、表格等）
    3. 保存原始 XML（保留完整格式）
    4. 如果有图片关系引用，保存 xml_rels 映射
    """
    type_name = "unknown"
    try:
        type_name = shape.shape_type.name if hasattr(shape, 'shape_type') else "unknown"
    except Exception:
        pass

    shape_info = {
        "type": type_name,
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
        "name": shape.name if hasattr(shape, 'name') else "",
    }

    # 占位符信息
    try:
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            ph_format = shape.placeholder_format
            shape_info["placeholder_idx"] = ph_format.idx
            shape_info["placeholder_type"] = ph_format.type.name
    except (ValueError, AttributeError):
        pass

    # 按形状类型提取内容
    has_content = False

    # 图片
    if _try_extract_image(shape, shape_info):
        has_content = True

    # 文本
    if _extract_text_content(shape, shape_info):
        has_content = True

    # 表格
    if type_name == "TABLE" or hasattr(shape, 'table'):
        try:
            _extract_table_data(shape, shape_info)
            has_content = True
        except Exception:
            pass

    # 组合形状
    if type_name == "GROUP_SHAPE" or hasattr(shape, 'shapes'):
        try:
            if hasattr(shape, 'shapes'):
                _extract_group_shapes(shape, shape_info)
                has_content = True
        except Exception:
            pass

    # 图表
    if type_name == "CHART":
        shape_info["is_chart"] = True
        try:
            if hasattr(shape, 'chart'):
                shape_info["chart_type"] = str(shape.chart.chart_type)
        except Exception:
            pass
        has_content = True

    # ⚠️ 关键：为所有形状保存原始 XML
    # 之前只对非 TEXT_BOX/PICTURE 保存，导致文本框丢失格式（填充、边框等）
    raw_xml = _extract_shape_xml(shape)
    if raw_xml:
        shape_info["raw_xml"] = raw_xml

        # 提取 XML 中的关系映射（图片引用等）
        if slide_part and all_images is not None:
            xml_rels = _extract_xml_relationships(shape, slide_part, all_images)
            if xml_rels:
                shape_info["xml_rels"] = xml_rels

    return shape_info


def extract_slide_to_dict(slide, slide_index: int, all_images: dict = None) -> Dict[str, Any]:
    """
    将单个幻灯片提取为字典。

    Args:
        slide: python-pptx Slide 对象
        slide_index: 幻灯片索引
        all_images: 可选的全局图片字典（用于 xml_rels 图片去重）
    """
    slide_info = {
        "index": slide_index,
        "shapes": [],
        "notes": "",
        "layout_name": "",
    }

    try:
        slide_info["layout_name"] = slide.slide_layout.name
    except Exception:
        pass

    # 获取 slide_part 以支持关系提取
    slide_part = slide.part if hasattr(slide, 'part') else None

    for shape in slide.shapes:
        try:
            shape_info = extract_shape_info(shape, slide_part, all_images)
            slide_info["shapes"].append(shape_info)
        except Exception as e:
            logger.warning(f"提取形状失败 (slide {slide_index}): {e}")
            fallback = {
                "type": "unknown",
                "left": 0, "top": 0, "width": 0, "height": 0,
                "name": "",
                "extraction_error": str(e),
            }
            try:
                fallback["name"] = shape.name
                fallback["left"] = shape.left
                fallback["top"] = shape.top
                fallback["width"] = shape.width
                fallback["height"] = shape.height
            except Exception:
                pass
            raw_xml = _extract_shape_xml(shape)
            if raw_xml:
                fallback["raw_xml"] = raw_xml
            slide_info["shapes"].append(fallback)

    # 备注
    try:
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            slide_info["notes"] = notes_slide.notes_text_frame.text
    except Exception:
        pass

    return slide_info


def extract_presentation(pptx_path: str) -> Dict[str, Any]:
    """
    提取整个演示文稿为字典。

    图片数据去重存储到顶层 images 字典。
    各形状保留 image_ref 引用。
    xml_rels 中的图片也存入 images 字典。
    """
    if not ensure_dependencies():
        return {"status": "error", "error": "依赖安装失败"}

    from pptx import Presentation

    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        return {"status": "error", "error": f"文件不存在: {pptx_path}"}

    try:
        prs = Presentation(str(pptx_path))

        pptx_data = {
            "source": str(pptx_path.name),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
            "slide_width_inches": prs.slide_width / 914400,
            "slide_height_inches": prs.slide_height / 914400,
            "slides": [],
            "layouts": [],
            "images": {},
        }

        # 提取布局信息
        for i, layout in enumerate(prs.slide_layouts):
            pptx_data["layouts"].append({
                "index": i,
                "name": layout.name
            })

        # 提取所有幻灯片
        # 传入 images 字典，以便 xml_rels 提取时图片也存入其中
        for i, slide in enumerate(prs.slides):
            slide_data = extract_slide_to_dict(slide, i, pptx_data["images"])

            # 将形状级别的图片数据移到统一的 images 字典
            for shape in slide_data["shapes"]:
                if "image_data" in shape:
                    img_hash = shape.get("image_hash", "")
                    if not img_hash:
                        img_hash = f"img_{i}_{shape.get('name', '0')}"
                        shape["image_hash"] = img_hash

                    if img_hash not in pptx_data["images"]:
                        pptx_data["images"][img_hash] = {
                            "type": shape.get("image_type", ""),
                            "data": shape["image_data"]
                        }
                    shape["image_ref"] = img_hash
                    del shape["image_data"]

            pptx_data["slides"].append(slide_data)

        return pptx_data

    except Exception as e:
        import traceback
        return {
            "status": "error",
            "error": f"提取失败: {str(e)}\n{traceback.format_exc()}"
        }
