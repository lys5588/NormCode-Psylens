#!/usr/bin/env python3
"""
PPTX 重建模块

提供两种重建策略：
1. build_from_template() — 幻灯片级克隆（最可靠）
   直接在 XML 和关系层面复制幻灯片，完美保留所有格式、图片引用等。
   适用于：应用模板计划、幻灯片选择/重排/复制。

2. rebuild_pptx() — 形状级重建（从 JSON 恢复）
   从提取的 JSON 数据重建幻灯片。
   适用于：文本替换后重建、跨格式转换。

关键修复：
- 使用 raw_xml 恢复所有形状类型（包括 TEXT_BOX），保留原始格式
- 清除布局自动生成的占位符，防止重复
- 关系映射：正确处理 raw_xml 中的图片引用（rId）
"""

import json
import base64
import logging
from copy import deepcopy
from io import BytesIO
from pathlib import Path
from typing import List, Dict, Any, Optional

logger = logging.getLogger(__name__)

# ============================================================
# 通用工具函数
# ============================================================

def _find_blank_layout(prs):
    """
    在演示文稿中查找空白布局。
    按优先级尝试多种匹配方式，而不是硬编码索引 6。
    """
    blank_patterns = ['blank', '空白', 'empty']
    for layout in prs.slide_layouts:
        name_lower = layout.name.lower()
        if any(pat in name_lower for pat in blank_patterns):
            return layout

    # 找占位符最少的布局
    min_ph = float('inf')
    best_layout = None
    for layout in prs.slide_layouts:
        try:
            n_ph = len(layout.placeholders)
            if n_ph < min_ph:
                min_ph = n_ph
                best_layout = layout
        except Exception:
            pass

    if best_layout is not None:
        return best_layout

    return prs.slide_layouts[-1] if len(prs.slide_layouts) > 0 else prs.slide_layouts[0]


def _find_layout_by_name(prs, layout_name: str):
    """按名称查找布局，未找到返回 None"""
    if not layout_name:
        return None
    for sl in prs.slide_layouts:
        if sl.name == layout_name:
            return sl
    return None


R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'

# 需要清除的形状元素标签
_SHAPE_TAGS = frozenset(('sp', 'pic', 'grpSp', 'graphicFrame', 'cxnSp'))

# 关系 ID 属性
_RID_ATTRS = [f'{{{R_NS}}}embed', f'{{{R_NS}}}link', f'{{{R_NS}}}id']


def _remap_relationship_ids(element, rId_map):
    """递归重映射 XML 元素中的所有关系 ID（r:embed, r:link, r:id）"""
    for attr in _RID_ATTRS:
        val = element.get(attr)
        if val and val in rId_map:
            element.set(attr, rId_map[val])
    for child in element:
        _remap_relationship_ids(child, rId_map)


def _clear_slide_shapes(slide):
    """清除幻灯片中由布局自动生成的占位符形状"""
    from lxml import etree
    sp_tree = slide.shapes._spTree
    to_remove = [child for child in sp_tree
                 if etree.QName(child.tag).localname in _SHAPE_TAGS]
    for elem in to_remove:
        sp_tree.remove(elem)
    return sp_tree


# ============================================================
# 策略一：幻灯片级克隆（最可靠）
# ============================================================

def _clone_slide_within_prs(prs, source_slide):
    """
    在同一演示文稿内克隆幻灯片。

    完整复制：
    - 所有形状（包括组合形状、图表、表格等）
    - 所有关系（图片、超链接等），正确重映射 rId
    - 幻灯片背景
    - 过渡效果和动画时间轴
    - 备注
    """
    from lxml import etree
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # 1. 清除布局自动生成的占位符
    sp_tree = _clear_slide_shapes(new_slide)

    # 2. 复制关系并构建 rId 映射
    rId_map = {}
    SKIP_RELTYPES = set()
    try:
        SKIP_RELTYPES.add(RT.SLIDE_LAYOUT)
    except Exception:
        pass
    try:
        SKIP_RELTYPES.add(RT.NOTES_SLIDE)
    except Exception:
        pass

    for rel in source_slide.part.rels.values():
        if rel.reltype in SKIP_RELTYPES:
            continue

        try:
            if rel.is_external:
                new_rId = new_slide.part.rels._next_rId
                new_slide.part.rels.add_relationship(
                    rel.reltype, rel.target_ref, new_rId, is_external=True
                )
                rId_map[rel.rId] = new_rId
            else:
                new_rId = new_slide.part.relate_to(rel.target_part, rel.reltype)
                rId_map[rel.rId] = new_rId
        except Exception as e:
            logger.warning(f"复制关系失败 {rel.rId}: {e}")

    # 3. 复制形状元素（重映射 rId）
    source_sp_tree = source_slide.shapes._spTree
    for child in source_sp_tree:
        tag = etree.QName(child.tag).localname
        if tag in _SHAPE_TAGS:
            new_elem = deepcopy(child)
            _remap_relationship_ids(new_elem, rId_map)
            sp_tree.append(new_elem)

    # 4. 复制幻灯片背景
    source_cSld = source_slide._element.find(f'{{{P_NS}}}cSld')
    target_cSld = new_slide._element.find(f'{{{P_NS}}}cSld')
    if source_cSld is not None and target_cSld is not None:
        source_bg = source_cSld.find(f'{{{P_NS}}}bg')
        if source_bg is not None:
            target_bg = target_cSld.find(f'{{{P_NS}}}bg')
            if target_bg is not None:
                target_cSld.remove(target_bg)
            new_bg = deepcopy(source_bg)
            _remap_relationship_ids(new_bg, rId_map)
            # bg 应在 cSld 的最前面（spTree 之前）
            target_cSld.insert(0, new_bg)

    # 5. 复制过渡效果和动画
    for tag_name in ('transition', 'timing'):
        source_elem = source_slide._element.find(f'{{{P_NS}}}{tag_name}')
        if source_elem is not None:
            target_elem = new_slide._element.find(f'{{{P_NS}}}{tag_name}')
            if target_elem is not None:
                new_slide._element.remove(target_elem)
            new_slide._element.append(deepcopy(source_elem))

    # 6. 复制备注
    try:
        if source_slide.has_notes_slide:
            notes_text = source_slide.notes_slide.notes_text_frame.text
            if notes_text.strip():
                new_slide.notes_slide.notes_text_frame.text = notes_text
    except Exception:
        pass

    return new_slide


def build_from_template(
    template_path: str,
    plan: List[int],
    output_path: str = None
) -> Dict[str, Any]:
    """
    通过幻灯片级克隆从模板构建新演示文稿。

    这是最可靠的方式——在 XML 和关系层面复制幻灯片，
    完美保留所有格式、图片、动画、背景等。

    工作原理：
    1. 打开模板
    2. 按计划将所需幻灯片克隆到末尾
    3. 删除前面的原始幻灯片
    4. 保存

    Args:
        template_path: 模板 PPTX 路径
        plan: 幻灯片索引列表（0-based，支持重复和重排）
        output_path: 输出路径

    Returns:
        dict: {status, output_path, slides_count, plan_applied}
    """
    try:
        from pptx_extract import ensure_dependencies
    except ImportError:
        from pptx_rebuild import _ensure_deps_fallback as ensure_dependencies

    if not ensure_dependencies():
        return {"status": "error", "error": "依赖安装失败"}

    from pptx import Presentation

    template_path = Path(template_path)
    if not template_path.exists():
        return {"status": "error", "error": f"模板不存在: {template_path}"}

    if not plan or not isinstance(plan, list):
        return {"status": "error", "error": "计划必须是非空的索引列表"}

    if output_path is None:
        output_path = template_path.parent / f"{template_path.stem}_planned.pptx"
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    try:
        prs = Presentation(str(template_path))
        original_count = len(prs.slides)

        # 索引越界回退：超出范围的索引使用最后一张幻灯片
        last_index = original_count - 1
        clamped_plan = [min(max(i, 0), last_index) for i in plan]
        if clamped_plan != plan:
            logger.warning(
                f"计划索引越界已回退: 原始={plan}, 修正={clamped_plan} "
                f"(模板共 {original_count} 张，索引 0-{last_index})"
            )
            plan = clamped_plan

        # 按计划克隆幻灯片（追加到末尾）
        for slide_index in plan:
            _clone_slide_within_prs(prs, prs.slides[slide_index])

        # 删除前面的原始幻灯片
        for _ in range(original_count):
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        prs.save(str(output_path))

        return {
            "status": "success",
            "output_path": str(output_path),
            "slides_count": len(prs.slides),
            "plan_applied": plan,
            "template_slides": original_count,
        }

    except Exception as e:
        import traceback
        return {
            "status": "error",
            "error": f"构建失败: {str(e)}\n{traceback.format_exc()}"
        }


# ============================================================
# 策略二：形状级重建（从 JSON 恢复）
# ============================================================

def _add_shape_from_xml(slide, raw_xml: str, xml_rels: dict = None,
                        pptx_data: dict = None) -> bool:
    """
    通过原始 XML 向幻灯片添加形状。

    关键改进：如果 xml_rels 提供了关系映射，
    会先恢复图片等资源的关系，然后重映射 rId 再插入 XML。
    这解决了"组合形状中的图片丢失导致文件损坏"的问题。

    Args:
        slide: 幻灯片对象
        raw_xml: 形状的原始 XML 字符串
        xml_rels: 可选的关系映射 {old_rId: {"image_ref": hash, "reltype": ...}}
        pptx_data: 可选的完整数据（用于查找 images）

    Returns:
        True if successful
    """
    try:
        from lxml import etree
        element = etree.fromstring(raw_xml)

        # 如果有关系映射，恢复关系并重映射 rId
        if xml_rels and pptx_data:
            rId_map = {}
            images = pptx_data.get("images", {})

            for old_rId, rel_info in xml_rels.items():
                img_ref = rel_info.get("image_ref", "")
                img_info = images.get(img_ref)

                if img_info and img_info.get("data"):
                    try:
                        img_bytes = base64.b64decode(img_info["data"])

                        # 使用 python-pptx 的正式 API 添加图片
                        # get_or_add_image_part 会正确注册到 Package
                        image_part, new_rId = slide.part.get_or_add_image_part(
                            BytesIO(img_bytes)
                        )
                        rId_map[old_rId] = new_rId
                    except Exception as e:
                        logger.warning(f"恢复关系 {old_rId} 失败: {e}")

            if rId_map:
                _remap_relationship_ids(element, rId_map)

        # 添加到形状树
        sp_tree = slide.shapes._spTree
        sp_tree.append(element)
        return True

    except Exception as e:
        logger.warning(f"通过 XML 添加形状失败: {e}")
        return False


def _add_image_shape(slide, shape_data: Dict, pptx_data: Dict) -> bool:
    """向幻灯片添加图片形状"""
    from pptx.util import Emu

    img_ref = shape_data.get("image_ref", "")
    img_data = pptx_data.get("images", {}).get(img_ref)

    if not img_data or not img_data.get("data"):
        if img_ref:
            logger.warning(f"找不到图片数据 (ref={img_ref})")
        return False

    try:
        img_bytes = base64.b64decode(img_data["data"])
        img_stream = BytesIO(img_bytes)
        left = Emu(shape_data.get("left", 0))
        top = Emu(shape_data.get("top", 0))
        width = Emu(shape_data.get("width", 914400))
        height = Emu(shape_data.get("height", 457200))
        slide.shapes.add_picture(img_stream, left, top, width, height)
        return True
    except Exception as e:
        logger.warning(f"添加图片失败 (ref={img_ref}): {e}")
        return False


def _add_text_shape(slide, shape_data: Dict) -> bool:
    """向幻灯片添加文本框形状（仅当无 raw_xml 时使用）"""
    from pptx.util import Emu, Pt

    paragraphs_data = shape_data.get("text_content", [])
    full_text = shape_data.get("full_text", "")

    if not paragraphs_data and not full_text:
        return False

    if not paragraphs_data and full_text:
        paragraphs_data = [{"text": full_text, "level": 0}]

    try:
        left = Emu(shape_data.get("left", 0))
        top = Emu(shape_data.get("top", 0))
        width = Emu(shape_data.get("width", 914400))
        height = Emu(shape_data.get("height", 457200))

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        for i, para_data in enumerate(paragraphs_data):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.level = para_data.get("level", 0)

            runs_data = para_data.get("runs", [])
            if runs_data:
                p.text = ""
                for j, run_data in enumerate(runs_data):
                    run = p.add_run()
                    run.text = run_data.get("text", "")
                    try:
                        if run_data.get("bold"):
                            run.font.bold = True
                        if run_data.get("italic"):
                            run.font.italic = True
                        if run_data.get("font_size"):
                            run.font.size = run_data["font_size"]
                        if run_data.get("font_name"):
                            run.font.name = run_data["font_name"]
                    except Exception:
                        pass
            else:
                p.text = para_data.get("text", "")

        return True
    except Exception as e:
        logger.warning(f"添加文本框失败 (name={shape_data.get('name', '?')}): {e}")
        return False


def add_shape_to_slide(slide, shape_data: Dict, pptx_data: Dict, prs=None) -> bool:
    """
    向幻灯片添加形状。

    改进策略：
    1. 所有有 raw_xml 的形状 → 优先用 XML 恢复（保留完整格式）
    2. 图片（无 raw_xml 的） → 用 API 添加
    3. 文本（无 raw_xml 的） → 用 API 添加

    Returns:
        True if shape was added, False if skipped/failed
    """
    shape_type = shape_data.get("type", "")
    has_raw_xml = "raw_xml" in shape_data
    has_image_ref = "image_ref" in shape_data
    has_text = bool(shape_data.get("text_content") or shape_data.get("full_text", "").strip())
    xml_rels = shape_data.get("xml_rels")

    # 优先用 raw_xml 恢复（保留所有格式：填充、边框、字体、阴影等）
    if has_raw_xml:
        if _add_shape_from_xml(slide, shape_data["raw_xml"], xml_rels, pptx_data):
            return True
        # XML 失败则尝试 API 方式

    # 图片（无 raw_xml 或 raw_xml 失败）
    if has_image_ref:
        return _add_image_shape(slide, shape_data, pptx_data)

    # 文本（无 raw_xml 或 raw_xml 失败）
    if has_text:
        return _add_text_shape(slide, shape_data)

    # 无法处理的形状
    logger.debug(f"跳过无法处理的形状: type={shape_type}, name={shape_data.get('name', '?')}")
    return False


def rebuild_pptx(
    pptx_data: Dict,
    output_path: str = None,
    slide_order: List[int] = None,
    template_path: str = None
) -> Dict[str, Any]:
    """
    从提取的数据重建 PPTX。

    Args:
        pptx_data: 由 pptx_extract.extract_presentation() 返回的数据字典
        output_path: 输出 PPTX 路径
        slide_order: 幻灯片顺序列表（支持重复和删除）
        template_path: 可选的模板路径（用作布局来源）

    Returns:
        dict: {status, output_path, slides_count}
    """
    try:
        from pptx_extract import ensure_dependencies
    except ImportError:
        from pptx_rebuild import _ensure_deps_fallback as ensure_dependencies

    if not ensure_dependencies():
        return {"status": "error", "error": "依赖安装失败"}

    from pptx import Presentation

    try:
        original_slides = pptx_data.get("slides", [])
        if slide_order is None:
            slide_order = list(range(len(original_slides)))

        # 创建新演示文稿
        if template_path and Path(template_path).exists():
            prs = Presentation(template_path)
            # 清空现有幻灯片
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
        else:
            prs = Presentation()

        # 设置幻灯片尺寸
        prs.slide_width = pptx_data.get("slide_width", 9144000)
        prs.slide_height = pptx_data.get("slide_height", 6858000)

        # 构建布局名称映射
        layout_map = {layout.name: layout for layout in prs.slide_layouts}
        blank_layout = _find_blank_layout(prs)

        shapes_added = 0
        shapes_skipped = 0

        for orig_index in slide_order:
            if orig_index < 0 or orig_index >= len(original_slides):
                logger.warning(f"跳过无效幻灯片索引: {orig_index}")
                continue

            slide_data = original_slides[orig_index]

            # 选择布局
            layout_name = slide_data.get("layout_name", "")
            layout = layout_map.get(layout_name, blank_layout)

            slide = prs.slides.add_slide(layout)

            # ⚠️ 关键修复：清除布局自动生成的占位符
            # 否则 "标题幻灯片" 等布局会添加额外的 Title/Subtitle 占位符
            _clear_slide_shapes(slide)

            # 添加形状
            for shape_data in slide_data.get("shapes", []):
                if add_shape_to_slide(slide, shape_data, pptx_data, prs):
                    shapes_added += 1
                else:
                    shapes_skipped += 1

            # 添加备注
            if slide_data.get("notes"):
                try:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = slide_data["notes"]
                except Exception:
                    pass

        # 保存
        if output_path is None:
            output_path = Path("rebuilt.pptx")
        else:
            output_path = Path(output_path)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        return {
            "status": "success",
            "output_path": str(output_path),
            "slides_count": len(prs.slides),
            "slide_order": slide_order,
            "shapes_added": shapes_added,
            "shapes_skipped": shapes_skipped,
        }

    except Exception as e:
        import traceback
        return {
            "status": "error",
            "error": f"重建失败: {str(e)}\n{traceback.format_exc()}"
        }


def rebuild_from_json(
    json_path: str,
    output_path: str = None,
    slide_order: List[int] = None,
    template_path: str = None
) -> Dict[str, Any]:
    """
    从 JSON 文件重建 PPTX（兼容旧接口）。
    """
    json_path = Path(json_path)
    if not json_path.exists():
        return {"status": "error", "error": f"JSON 文件不存在: {json_path}"}

    try:
        with open(json_path, 'r', encoding='utf-8') as f:
            pptx_data = json.load(f)
    except Exception as e:
        return {"status": "error", "error": f"读取 JSON 失败: {e}"}

    if output_path is None:
        output_path = str(json_path.parent / f"{json_path.stem}_rebuilt.pptx")

    return rebuild_pptx(pptx_data, output_path, slide_order, template_path)


def manipulate_and_rebuild(
    json_path: str,
    operations: List[Dict],
    output_path: str = None
) -> Dict[str, Any]:
    """
    对幻灯片执行操作（复制、删除、重排）后重建。
    """
    json_path = Path(json_path)
    if not json_path.exists():
        return {"status": "error", "error": f"JSON 文件不存在: {json_path}"}

    try:
        with open(json_path, 'r', encoding='utf-8') as f:
            pptx_data = json.load(f)

        n_slides = len(pptx_data.get("slides", []))
        order = list(range(n_slides))

        for op in operations:
            action = op.get("action", "")

            if action == "duplicate":
                idx = op.get("index", 0)
                if 0 <= idx < len(order):
                    order.insert(idx + 1, order[idx])

            elif action == "delete":
                idx = op.get("index", 0)
                if 0 <= idx < len(order):
                    order.pop(idx)

            elif action == "move":
                from_idx = op.get("from", 0)
                to_idx = op.get("to", 0)
                if 0 <= from_idx < len(order):
                    item = order.pop(from_idx)
                    to_idx = min(to_idx, len(order))
                    order.insert(to_idx, item)

            elif action == "reorder":
                new_order = op.get("order", [])
                if new_order:
                    order = [i for i in new_order if 0 <= i < n_slides]

        return rebuild_pptx(pptx_data, output_path, order)

    except Exception as e:
        import traceback
        return {
            "status": "error",
            "error": f"操作失败: {str(e)}\n{traceback.format_exc()}"
        }


def _ensure_deps_fallback() -> bool:
    """Fallback dependency check when pptx_extract is not importable"""
    try:
        from pptx import Presentation
        from lxml import etree
        return True
    except ImportError:
        return False
