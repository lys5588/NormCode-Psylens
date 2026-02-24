#!/usr/bin/env python3
"""
优化文本布局脚本

分析 PPTX 中的文本框，检测溢出问题并自动调整字体大小。

输入：{演示文稿路径}, {输出路径}（可选）
输出：{status, output_path, adjustments_made, report}
"""

import json
import subprocess
import sys
from pathlib import Path
from typing import Dict, List, Any, Optional


def ensure_pptx_installed():
    """确保 python-pptx 已安装"""
    try:
        import pptx
        return True
    except ImportError:
        try:
            subprocess.check_call([
                sys.executable, "-m", "pip", "install", 
                "python-pptx>=0.6.21", "--quiet"
            ])
            return True
        except Exception:
            return False


def estimate_text_width(text: str, font_size_pt: float) -> float:
    """
    估算文本宽度（EMU）
    
    使用经验公式：中文字符约等于字体大小，英文约等于0.5倍字体大小
    """
    chinese_chars = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
    other_chars = len(text) - chinese_chars
    
    # 中文字符宽度约等于字体高度，英文约0.5倍
    estimated_width_pt = chinese_chars * font_size_pt + other_chars * font_size_pt * 0.5
    
    # 转换为 EMU (1 pt = 12700 EMU)
    return estimated_width_pt * 12700


def analyze_text_frame(shape, slide_index: int, shape_index: int) -> Optional[Dict]:
    """
    分析单个文本框，检测是否需要调整
    
    Returns:
        调整建议或 None
    """
    if not hasattr(shape, 'text_frame') or not shape.has_text_frame:
        return None
    
    tf = shape.text_frame
    text = tf.text.strip()
    
    if not text:
        return None
    
    # 获取形状尺寸
    shape_width = shape.width  # EMU
    shape_height = shape.height  # EMU
    
    # 获取字体大小
    font_size = None
    if len(tf.paragraphs) > 0 and len(tf.paragraphs[0].runs) > 0:
        font_size = tf.paragraphs[0].runs[0].font.size
    
    if not font_size:
        return None
    
    font_size_pt = font_size / 12700  # EMU to pt
    
    # 估算文本宽度
    estimated_width = estimate_text_width(text, font_size_pt)
    
    # 检测溢出（留10%边距）
    available_width = shape_width * 0.9
    
    if estimated_width > available_width:
        # 计算建议的字体大小
        ratio = available_width / estimated_width
        suggested_size_pt = max(8, font_size_pt * ratio)  # 最小8pt
        
        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "shape_name": shape.name,
            "text_preview": text[:50] + "..." if len(text) > 50 else text,
            "current_font_pt": round(font_size_pt, 1),
            "suggested_font_pt": round(suggested_size_pt, 1),
            "overflow_ratio": round(estimated_width / available_width, 2),
            "action": "reduce_font"
        }
    
    return None


def adjust_font_size(shape, new_size_pt: float) -> bool:
    """调整形状中所有文本的字体大小"""
    from pptx.util import Pt
    
    if not hasattr(shape, 'text_frame') or not shape.has_text_frame:
        return False
    
    try:
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                run.font.size = Pt(new_size_pt)
        return True
    except Exception:
        return False


def optimize_presentation(
    pptx_path: str,
    output_path: str = None,
    auto_adjust: bool = True,
    min_font_pt: float = 8,
    body=None
) -> Dict[str, Any]:
    """
    分析并优化演示文稿中的文本布局
    
    Args:
        pptx_path: 输入 PPTX 路径
        output_path: 输出路径（可选，默认覆盖原文件）
        auto_adjust: 是否自动调整字体大小
        min_font_pt: 最小字体大小
        body: Agent body
    
    Returns:
        dict: {status, output_path, adjustments_made, report}
    """
    result = {
        "status": "pending",
        "output_path": "",
        "adjustments_made": 0,
        "report": []
    }
    
    if not ensure_pptx_installed():
        result["status"] = "error"
        result["error"] = "无法安装 python-pptx"
        return result
    
    try:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.enum.text import MSO_AUTO_SIZE
        
        # 通过 file_system.resolve_path 解析路径
        if body and hasattr(body, 'file_system') and hasattr(body.file_system, 'resolve_path'):
            resolved_path = body.file_system.resolve_path(pptx_path)
        else:
            resolved_path = str(Path(pptx_path).resolve())
        
        if not Path(resolved_path).exists():
            result["status"] = "error"
            result["error"] = f"文件不存在: {pptx_path} (resolved: {resolved_path})"
            return result
        
        # 确定输出路径（通过 file_system.resolve_path 解析）
        if output_path:
            if body and hasattr(body, 'file_system') and hasattr(body.file_system, 'resolve_path'):
                output_path = Path(body.file_system.resolve_path(str(output_path)))
            else:
                output_path = Path(output_path)
        else:
            # 默认生成 _optimized 后缀
            p = Path(resolved_path)
            output_path = p.parent / f"{p.stem}_optimized.pptx"
        
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # 打开演示文稿
        prs = Presentation(resolved_path)
        
        adjustments = []
        adjustments_made = 0
        
        # 遍历所有幻灯片
        for slide_idx, slide in enumerate(prs.slides):
            shapes_list = list(slide.shapes)
            
            for shape_idx, shape in enumerate(shapes_list):
                # 分析文本框
                suggestion = analyze_text_frame(shape, slide_idx, shape_idx)
                
                if suggestion:
                    adjustments.append(suggestion)
                    
                    # 如果启用自动调整
                    if auto_adjust:
                        suggested_size = max(min_font_pt, suggestion["suggested_font_pt"])
                        if adjust_font_size(shape, suggested_size):
                            adjustments_made += 1
                            suggestion["adjusted"] = True
                            suggestion["final_font_pt"] = suggested_size
                        else:
                            suggestion["adjusted"] = False
                    
                    # 尝试启用自动收缩
                    try:
                        if hasattr(shape, 'text_frame'):
                            shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    except Exception:
                        pass
        
        # 保存
        prs.save(str(output_path))
        
        result["status"] = "success"
        result["output_path"] = str(output_path)
        result["adjustments_made"] = adjustments_made
        result["report"] = adjustments
        result["total_issues_found"] = len(adjustments)
        
        # 生成摘要
        if adjustments:
            result["summary"] = f"发现 {len(adjustments)} 处潜在溢出，已自动调整 {adjustments_made} 处"
        else:
            result["summary"] = "未发现文本溢出问题"
        
        return result
        
    except Exception as e:
        import traceback
        result["status"] = "error"
        result["error"] = f"优化失败: {str(e)}"
        result["traceback"] = traceback.format_exc()
        return result


def flatten_input(data):
    """展平嵌套的输入"""
    if isinstance(data, list):
        while isinstance(data, list) and len(data) > 0:
            if isinstance(data[0], (dict, str)):
                break
            data = data[0]
    return data


def main(
    input_1=None,  # 演示文稿路径
    input_2=None,  # 输出路径（可选）
    body=None,
    **kwargs
) -> Dict:
    """
    主函数 - 优化文本布局
    
    Args:
        input_1: {演示文稿路径} 或 {基础演示文稿} 的结果
        input_2: {输出路径} - 可选
        body: Agent body
    
    Returns:
        dict: 处理结果
    """
    try:
        # 解析输入路径
        pptx_input = flatten_input(input_1)
        
        # 如果输入是字典（如前一步的输出结果），提取路径
        if isinstance(pptx_input, dict):
            pptx_path = pptx_input.get("output_path") or pptx_input.get("path", "")
        else:
            pptx_path = str(pptx_input) if pptx_input else ""
        
        # 解析输出路径
        output_path = flatten_input(input_2)
        if isinstance(output_path, dict):
            output_path = output_path.get("path", None)
        if output_path and not isinstance(output_path, str):
            output_path = str(output_path)
        
        return optimize_presentation(
            pptx_path=pptx_path,
            output_path=output_path,
            auto_adjust=True,
            body=body
        )
        
    except Exception as e:
        import traceback
        return {
            "status": "error",
            "error": f"处理失败: {str(e)}",
            "traceback": traceback.format_exc()
        }


if __name__ == "__main__":
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')
    
    if len(sys.argv) >= 2:
        pptx_file = sys.argv[1]
        output = sys.argv[2] if len(sys.argv) > 2 else None
        
        result = main(pptx_file, output)
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        print("""
优化文本布局工具

用法:
    python 优化文本布局.py <演示文稿> [输出路径]

参数:
    演示文稿: PPTX 文件路径
    输出路径: 可选，默认生成 xxx_optimized.pptx

功能:
    1. 分析所有文本框，检测潜在溢出
    2. 自动调整字体大小以适应容器
    3. 启用文本自动收缩
    4. 生成调整报告

示例:
    python 优化文本布局.py presentation.pptx
    python 优化文本布局.py input.pptx output.pptx
""")

