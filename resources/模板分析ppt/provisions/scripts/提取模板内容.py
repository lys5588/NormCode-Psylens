#!/usr/bin/env python3
"""
提取模板内容脚本

从 PPTX 模板中提取所有页面的原始数据，包括：
- 每个形状的类型、位置、尺寸
- 文本内容和字体样式
- 占位符信息
- 图片信息

⚠️ 二进制约束：读取 PPTX + 提取数据必须一步完成
（PPTX 对象不可序列化，无法跨 step 传递）

输入：{模板路径}
输出：{模板原始数据} - 包含每个页面的完整原始数据
"""

import json
import subprocess
import sys
from pathlib import Path
from typing import Dict, List, Any, Optional


def ensure_pptx_installed():
    """确保 python-pptx 已安装"""
    try:
        import pptx
        return True
    except ImportError:
        try:
            subprocess.check_call([
                sys.executable, "-m", "pip", "install", 
                "python-pptx>=0.6.21", "--quiet"
            ])
            import importlib
            import site
            importlib.invalidate_caches()
            if hasattr(site, 'main'):
                site.main()
            return True
        except Exception:
            return False


def emu_to_cm(emu_value) -> float:
    """将EMU转换为厘米"""
    if emu_value is None:
        return 0
    return round(emu_value / 360000, 2)


def emu_to_pt(emu_value) -> float:
    """将EMU转换为点"""
    if emu_value is None:
        return 0
    return round(emu_value / 12700, 1)


def extract_font_info(run) -> Dict:
    """提取字体信息"""
    font = run.font
    info = {}
    
    try:
        if font.size:
            info["size_pt"] = font.size.pt
    except Exception:
        pass
    
    try:
        if font.bold is not None:
            info["bold"] = font.bold
    except Exception:
        pass
    
    try:
        if font.italic is not None:
            info["italic"] = font.italic
    except Exception:
        pass
    
    try:
        if font.name:
            info["name"] = font.name
    except Exception:
        pass
    
    # 颜色提取需要特殊处理 - scheme colors 没有 .rgb 属性
    try:
        if font.color:
            color_type = font.color.type
            if color_type is not None:
                # 尝试获取 RGB 值
                try:
                    if hasattr(font.color, 'rgb') and font.color.rgb:
                        info["color"] = str(font.color.rgb)
                except (AttributeError, TypeError):
                    # 对于 scheme colors，尝试获取主题颜色名称
                    try:
                        if hasattr(font.color, 'theme_color') and font.color.theme_color:
                            info["theme_color"] = str(font.color.theme_color)
                    except Exception:
                        pass
    except Exception:
        pass
    
    return info if info else None


def extract_paragraph_info(para) -> Dict:
    """提取段落信息"""
    info = {
        "text": para.text,
        "level": para.level,
    }
    
    # 段落对齐
    try:
        if para.alignment:
            info["alignment"] = str(para.alignment)
    except Exception:
        pass
    
    # 提取第一个run的字体信息作为段落样式
    try:
        for run in para.runs:
            if run.text.strip():
                font_info = extract_font_info(run)
                if font_info:
                    info["font"] = font_info
                break
    except Exception:
        pass
    
    return info


def extract_text_frame(shape) -> Optional[Dict]:
    """提取文本框完整信息"""
    try:
        if not hasattr(shape, 'text_frame') or not shape.has_text_frame:
            return None
        
        tf = shape.text_frame
        paragraphs = []
        
        for para in tf.paragraphs:
            try:
                para_info = extract_paragraph_info(para)
                paragraphs.append(para_info)
            except Exception:
                # 即使单个段落失败，也继续处理其他段落
                paragraphs.append({"text": para.text if hasattr(para, 'text') else "", "level": 0})
        
        return {
            "full_text": tf.text,
            "paragraphs": paragraphs
        }
    except Exception:
        # 尝试至少获取文本
        try:
            if hasattr(shape, 'text'):
                return {"full_text": shape.text, "paragraphs": []}
        except Exception:
            pass
        return None


def extract_shape(shape, shape_idx: int, slide_width, slide_height) -> Dict:
    """提取单个形状的原始数据"""
    # 基础位置信息 - 使用 try/except 保护每个属性
    position = {}
    try:
        position["left_cm"] = emu_to_cm(shape.left)
        position["top_cm"] = emu_to_cm(shape.top)
        position["width_cm"] = emu_to_cm(shape.width)
        position["height_cm"] = emu_to_cm(shape.height)
        position["left_emu"] = shape.left
        position["top_emu"] = shape.top
        position["width_emu"] = shape.width
        position["height_emu"] = shape.height
    except Exception:
        pass
    
    # 基础信息
    data = {
        "shape_index": shape_idx,
        "element_type": "other",  # 默认类型
    }
    
    try:
        data["shape_id"] = shape.shape_id if hasattr(shape, 'shape_id') else shape_idx
    except Exception:
        data["shape_id"] = shape_idx
    
    try:
        data["name"] = shape.name if hasattr(shape, 'name') else ""
    except Exception:
        data["name"] = ""
    
    try:
        data["shape_type"] = str(shape.shape_type) if hasattr(shape, 'shape_type') else "Unknown"
    except Exception:
        data["shape_type"] = "Unknown"
    
    data["position"] = position
    
    # 占位符信息
    try:
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            data["is_placeholder"] = True
            data["placeholder_info"] = {
                "idx": shape.placeholder_format.idx,
                "type": str(shape.placeholder_format.type)
            }
        else:
            data["is_placeholder"] = False
    except Exception:
        data["is_placeholder"] = False
    
    # 判断形状类型并提取相应内容
    shape_type_str = data.get("shape_type", "")
    
    # 图片
    if "PICTURE" in shape_type_str:
        data["element_type"] = "picture"
        # 计算图片覆盖率（用于后续判断是否为背景）
        try:
            if slide_width and slide_height and position.get("width_emu") and position.get("height_emu"):
                shape_area = position["width_emu"] * position["height_emu"]
                slide_area = slide_width * slide_height
                data["coverage_ratio"] = round(shape_area / slide_area, 3) if slide_area > 0 else 0
        except Exception:
            pass
        try:
            if hasattr(shape, 'image'):
                data["image_info"] = {
                    "content_type": shape.image.content_type,
                    "size_bytes": len(shape.image.blob) if shape.image.blob else 0
                }
        except Exception:
            pass
        return data
    
    # 表格
    if "TABLE" in shape_type_str:
        data["element_type"] = "table"
        try:
            if hasattr(shape, 'table'):
                table = shape.table
                data["table_info"] = {
                    "rows": len(table.rows),
                    "columns": len(table.columns)
                }
                # 提取表格内容
                cells = []
                for row_idx, row in enumerate(table.rows):
                    for col_idx, cell in enumerate(row.cells):
                        cells.append({
                            "row": row_idx,
                            "col": col_idx,
                            "text": cell.text
                        })
                data["table_cells"] = cells
        except Exception:
            pass
        return data
    
    # 图表
    if "CHART" in shape_type_str:
        data["element_type"] = "chart"
        try:
            if hasattr(shape, 'chart'):
                data["chart_info"] = {
                    "chart_type": str(shape.chart.chart_type) if hasattr(shape.chart, 'chart_type') else "Unknown"
                }
        except Exception:
            pass
        return data
    
    # 文本框/形状 - 优先尝试提取文本
    try:
        text_frame = extract_text_frame(shape)
        if text_frame:
            full_text = text_frame.get("full_text", "")
            if full_text.strip():
                data["element_type"] = "text"
                data["text_content"] = full_text
                data["text_frame"] = text_frame
                return data
            else:
                # 有文本框但没有文本 - 可能是占位符
                data["text_frame"] = text_frame
    except Exception:
        pass
    
    # 检查 TEXT_BOX 类型
    if "TEXT_BOX" in shape_type_str or "TEXT" in shape_type_str:
        data["element_type"] = "text"
        # 尝试直接获取文本
        try:
            if hasattr(shape, 'text') and shape.text:
                data["text_content"] = shape.text
        except Exception:
            pass
        return data
    
    # 其他形状类型
    if "AUTO_SHAPE" in shape_type_str or "SHAPE" in shape_type_str:
        data["element_type"] = "shape"
    elif "GROUP" in shape_type_str:
        data["element_type"] = "group"
    else:
        data["element_type"] = "other"
    
    return data


def extract_page(slide, page_idx: int, total_pages: int, prs) -> Dict:
    """提取单个页面的原始数据"""
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    page_data = {
        "index": page_idx,
        "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
        "shapes": []
    }
    
    # 提取每个形状
    for shape_idx, shape in enumerate(slide.shapes):
        try:
            shape_data = extract_shape(shape, shape_idx, slide_width, slide_height)
            page_data["shapes"].append(shape_data)
        except Exception as e:
            # 即使提取失败，也尝试获取基本信息和文本
            fallback_data = {
                "shape_index": shape_idx,
                "element_type": "other",
                "extraction_error": str(e)
            }
            
            # 尝试获取形状名称
            try:
                if hasattr(shape, 'name'):
                    fallback_data["name"] = shape.name
            except Exception:
                pass
            
            # 尝试获取形状类型
            try:
                if hasattr(shape, 'shape_type'):
                    fallback_data["shape_type"] = str(shape.shape_type)
            except Exception:
                pass
            
            # ⚠️ 关键：尝试直接获取文本内容
            try:
                if hasattr(shape, 'text_frame') and shape.has_text_frame:
                    text = shape.text_frame.text
                    if text and text.strip():
                        fallback_data["element_type"] = "text"
                        fallback_data["text_content"] = text
            except Exception:
                pass
            
            # 如果上面失败，尝试直接访问 .text 属性
            if fallback_data.get("element_type") != "text":
                try:
                    if hasattr(shape, 'text'):
                        text = shape.text
                        if text and text.strip():
                            fallback_data["element_type"] = "text"
                            fallback_data["text_content"] = text
                except Exception:
                    pass
            
            page_data["shapes"].append(fallback_data)
    
    # 页面统计信息
    page_data["stats"] = {
        "total_shapes": len(page_data["shapes"]),
        "text_shapes": sum(1 for s in page_data["shapes"] if s.get("element_type") == "text"),
        "picture_shapes": sum(1 for s in page_data["shapes"] if s.get("element_type") == "picture"),
        "placeholder_count": sum(1 for s in page_data["shapes"] if s.get("is_placeholder")),
    }
    
    return page_data


def extract_template_content(template_path: str, body=None) -> Dict:
    """
    提取 PPT 模板的所有页面内容
    
    Args:
        template_path: PPT模板文件路径
        body: Agent body (可选)
    
    Returns:
        模板原始数据
    """
    if not ensure_pptx_installed():
        return {
            "status": "error",
            "error": "无法安装 python-pptx 库",
            "template_path": template_path
        }
    
    try:
        from pptx import Presentation
        
        # 通过 file_system.resolve_path 解析模板路径
        if body and hasattr(body, 'file_system') and hasattr(body.file_system, 'resolve_path'):
            resolved_path = body.file_system.resolve_path(template_path)
        else:
            resolved_path = str(Path(template_path).resolve())
        
        if not Path(resolved_path).exists():
            return {
                "status": "error",
                "error": f"无法找到模板文件: {template_path} (resolved: {resolved_path})",
                "template_path": template_path
            }
        
        # ⚠️ 关键：读取 PPTX + 提取数据在同一步完成
        prs = Presentation(resolved_path)
        total_pages = len(prs.slides)
        
        # 提取每个页面
        pages = []
        for i, slide in enumerate(prs.slides):
            page_data = extract_page(slide, i, total_pages, prs)
            pages.append(page_data)
        
        # 返回原始数据
        return {
            "status": "success",
            "template_path": resolved_path,
            "template_info": {
                "total_pages": total_pages,
                "slide_width_cm": emu_to_cm(prs.slide_width),
                "slide_height_cm": emu_to_cm(prs.slide_height),
                "slide_width_emu": prs.slide_width,
                "slide_height_emu": prs.slide_height
            },
            "pages": pages
        }
        
    except Exception as e:
        import traceback
        return {
            "status": "error",
            "error": f"提取模板内容失败: {str(e)}",
            "traceback": traceback.format_exc(),
            "template_path": template_path
        }


def flatten_input(data):
    """展平嵌套的输入"""
    if isinstance(data, list):
        while isinstance(data, list) and len(data) > 0:
            data = data[0]
    return data


def main(input_1=None, body=None, **kwargs) -> Dict:
    """
    主函数 - 提取模板内容
    
    Args:
        input_1: {模板路径}
        body: Agent body
    
    Returns:
        dict: 模板原始数据
    """
    try:
        template_path = flatten_input(input_1)
        if not isinstance(template_path, str):
            template_path = str(template_path) if template_path else ""
        
        return extract_template_content(template_path, body)
        
    except Exception as e:
        return {
            "status": "error",
            "error": f"提取模板内容失败: {str(e)}",
            "template_path": str(input_1) if input_1 else ""
        }


if __name__ == "__main__":
    import sys
    
    if len(sys.argv) >= 2:
        template_path = sys.argv[1]
        result = main(input_1=template_path)
        print(json.dumps(result, ensure_ascii=False, indent=2, default=str))
    else:
        print(json.dumps({
            "error": "用法: python 提取模板内容.py <template_path>"
        }, ensure_ascii=False))

