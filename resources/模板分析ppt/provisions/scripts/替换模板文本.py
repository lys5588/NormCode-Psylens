#!/usr/bin/env python3
"""
替换模板文本脚本

根据替换计划，在现有模板页面上替换指定的文本内容。
不创建新页面，只修改现有页面的文本。

⚠️ 二进制约束：读取 PPTX + 修改 + 保存必须一步完成

输入：
  - {模板路径}: PPTX 模板文件
  - {替换计划}: 指定哪些文本需要替换
  - {输出路径}: 保存修改后的文件（可选）

替换计划格式：
[
  {
    "page_index": 0,
    "replacements": [
      {"shape_index": 10, "new_text": "新标题"},
      {"shape_index": 11, "new_text": "新副标题"}
    ]
  },
  {
    "page_index": 1,
    "replacements": [
      {"shape_index": 5, "new_text": "第二页内容"}
    ]
  }
]

输出：{status, output_path, changes_made}
"""

import json
import subprocess
import sys
from pathlib import Path
from typing import Dict, List, Any, Optional


def ensure_pptx_installed():
    """确保 python-pptx 已安装"""
    try:
        import pptx
        return True
    except ImportError:
        try:
            subprocess.check_call([
                sys.executable, "-m", "pip", "install", 
                "python-pptx>=0.6.21", "--quiet"
            ])
            import importlib
            import site
            importlib.invalidate_caches()
            if hasattr(site, 'main'):
                site.main()
            return True
        except Exception:
            return False


def replace_shape_text(
    shape, 
    new_text: str, 
    preserve_formatting: bool = True,
    auto_fit: bool = True,
    adjust_font_size: bool = False,
    min_font_size_pt: int = 8
) -> bool:
    """
    替换形状中的文本，支持自动调整
    
    Args:
        shape: pptx shape 对象
        new_text: 新的文本内容
        preserve_formatting: 是否保留原有格式
        auto_fit: 是否启用自动收缩文本以适应形状
        adjust_font_size: 是否根据文本长度智能调整字体大小
        min_font_size_pt: 最小字体大小（磅）
    
    Returns:
        bool: 是否成功替换
    """
    if not hasattr(shape, 'text_frame') or not shape.has_text_frame:
        return False
    
    tf = shape.text_frame
    
    # 启用自动换行
    tf.word_wrap = True
    
    # 启用自动收缩文本以适应形状
    if auto_fit:
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        except Exception:
            pass  # 某些形状可能不支持
    
    # 获取原始字体大小（用于智能调整）
    original_font_size = None
    if len(tf.paragraphs) > 0 and len(tf.paragraphs[0].runs) > 0:
        original_font_size = tf.paragraphs[0].runs[0].font.size
    
    if preserve_formatting and len(tf.paragraphs) > 0:
        # 保留格式：只替换第一个段落的文本，保持样式
        # 对于多行文本，用换行符分割
        lines = new_text.split('\n')
        
        # 清空现有段落
        for i, para in enumerate(tf.paragraphs):
            if i < len(lines):
                # 保留第一个 run 的格式
                if len(para.runs) > 0:
                    first_run = para.runs[0]
                    # 清空其他 runs
                    for run in para.runs[1:]:
                        run.text = ""
                    first_run.text = lines[i]
                    
                    # 智能调整字体大小
                    if adjust_font_size and original_font_size:
                        adjusted_size = calculate_adjusted_font_size(
                            lines[i], 
                            original_font_size,
                            shape.width,
                            min_font_size_pt
                        )
                        if adjusted_size:
                            first_run.font.size = adjusted_size
                else:
                    para.text = lines[i]
            else:
                # 多余的段落清空
                para.text = ""
        
        # 如果新文本行数更多，在最后一个段落追加
        if len(lines) > len(tf.paragraphs):
            remaining = '\n'.join(lines[len(tf.paragraphs):])
            if len(tf.paragraphs) > 0:
                last_para = tf.paragraphs[-1]
                if len(last_para.runs) > 0:
                    last_para.runs[0].text += '\n' + remaining
                else:
                    last_para.text += '\n' + remaining
    else:
        # 不保留格式：直接替换
        tf.text = new_text
    
    return True


def calculate_adjusted_font_size(text: str, original_size, shape_width, min_size_pt: int = 8):
    """
    根据文本长度和形状宽度计算调整后的字体大小
    
    Args:
        text: 文本内容
        original_size: 原始字体大小 (EMU)
        shape_width: 形状宽度 (EMU)
        min_size_pt: 最小字体大小（磅）
    
    Returns:
        调整后的字体大小 (EMU) 或 None
    """
    from pptx.util import Pt, Emu
    
    if not original_size or not shape_width:
        return None
    
    # 估算每个字符的平均宽度（粗略估计：字体大小 * 0.6）
    original_pt = original_size / 12700  # EMU to Pt
    estimated_char_width = original_pt * 0.6 * 12700  # 每字符宽度 in EMU
    
    # 计算文本需要的宽度
    text_width = len(text) * estimated_char_width
    
    # 如果文本太长，需要缩小字体
    if text_width > shape_width * 0.9:  # 留10%边距
        ratio = (shape_width * 0.9) / text_width
        new_pt = original_pt * ratio
        
        # 确保不小于最小字体
        if new_pt < min_size_pt:
            new_pt = min_size_pt
        
        return Pt(new_pt)
    
    return None  # 不需要调整


def substitute_template_text(
    template_path: str,
    substitution_plan: List[Dict],
    output_path: str = None,
    preserve_formatting: bool = True,
    auto_fit: bool = True,
    adjust_font_size: bool = False,
    body=None
) -> Dict[str, Any]:
    """
    根据替换计划替换模板中的文本
    
    Args:
        template_path: 模板 PPTX 路径
        substitution_plan: 替换计划列表
        output_path: 输出路径（可选）
        preserve_formatting: 是否保留原有格式
        body: Agent body（可选）
    
    Returns:
        dict: {status, output_path, changes_made, details}
    """
    result = {
        "status": "pending",
        "output_path": "",
        "changes_made": 0,
        "details": []
    }
    
    if not ensure_pptx_installed():
        result["status"] = "error"
        result["error"] = "无法安装 python-pptx 库"
        return result
    
    try:
        from pptx import Presentation
        
        # 通过 file_system.resolve_path 解析模板路径
        if body and hasattr(body, 'file_system') and hasattr(body.file_system, 'resolve_path'):
            resolved_path = body.file_system.resolve_path(template_path)
        else:
            resolved_path = str(Path(template_path).resolve())
        
        if not Path(resolved_path).exists():
            result["status"] = "error"
            result["error"] = f"无法找到模板文件: {template_path} (resolved: {resolved_path})"
            return result
        
        # 确定输出路径（通过 file_system.resolve_path 解析）
        if output_path:
            if body and hasattr(body, 'file_system') and hasattr(body.file_system, 'resolve_path'):
                output_path = Path(body.file_system.resolve_path(str(output_path)))
            else:
                output_path = Path(output_path)
        else:
            template_p = Path(resolved_path)
            output_path = template_p.parent / f"{template_p.stem}_substituted.pptx"
        
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # ⚠️ 关键：读取 PPTX + 修改 + 保存在同一步完成
        prs = Presentation(resolved_path)
        total_slides = len(prs.slides)
        
        changes_made = 0
        details = []
        
        # 按页面处理替换
        for page_plan in substitution_plan:
            # 支持 page_index 或 slide_index（兼容不同命名）
            page_index = page_plan.get("page_index") or page_plan.get("slide_index")
            replacements = page_plan.get("replacements", [])
            
            if page_index is None or page_index < 0 or page_index >= total_slides:
                details.append({
                    "page_index": page_index,
                    "status": "skipped",
                    "reason": f"无效的页面索引（共 {total_slides} 页）"
                })
                continue
            
            slide = prs.slides[page_index]
            page_changes = []
            
            for replacement in replacements:
                shape_index = replacement.get("shape_index")
                shape_id = replacement.get("shape_id")
                # 支持 new_text 或 text（兼容不同命名）
                new_text = replacement.get("new_text") or replacement.get("text", "")
                
                # 查找目标 shape
                target_shape = None
                
                if shape_index is not None:
                    # 按索引查找
                    shapes_list = list(slide.shapes)
                    if 0 <= shape_index < len(shapes_list):
                        target_shape = shapes_list[shape_index]
                
                if target_shape is None and shape_id is not None:
                    # 按 shape_id 查找
                    for shape in slide.shapes:
                        if hasattr(shape, 'shape_id') and shape.shape_id == shape_id:
                            target_shape = shape
                            break
                
                if target_shape is None:
                    page_changes.append({
                        "shape_index": shape_index,
                        "shape_id": shape_id,
                        "status": "not_found"
                    })
                    continue
                
                # 获取原文本
                old_text = ""
                if hasattr(target_shape, 'text_frame') and target_shape.has_text_frame:
                    old_text = target_shape.text_frame.text
                
                # 替换文本（启用自动适应）
                success = replace_shape_text(
                    target_shape, 
                    new_text, 
                    preserve_formatting=preserve_formatting,
                    auto_fit=auto_fit,
                    adjust_font_size=adjust_font_size
                )
                
                if success:
                    changes_made += 1
                    page_changes.append({
                        "shape_index": shape_index,
                        "shape_id": shape_id,
                        "shape_name": target_shape.name,
                        "old_text": old_text[:100] + "..." if len(old_text) > 100 else old_text,
                        "new_text": new_text[:100] + "..." if len(new_text) > 100 else new_text,
                        "status": "replaced"
                    })
                else:
                    page_changes.append({
                        "shape_index": shape_index,
                        "shape_id": shape_id,
                        "status": "no_text_frame"
                    })
            
            details.append({
                "page_index": page_index,
                "changes": page_changes
            })
        
        # 保存修改后的文件
        prs.save(str(output_path))
        
        result["status"] = "success"
        result["output_path"] = str(output_path)
        result["changes_made"] = changes_made
        result["details"] = details
        result["total_slides"] = total_slides
        
        return result
        
    except Exception as e:
        import traceback
        result["status"] = "error"
        result["error"] = f"替换文本失败: {str(e)}"
        result["traceback"] = traceback.format_exc()
        return result


def flatten_input(data):
    """展平嵌套的输入"""
    if isinstance(data, list):
        while isinstance(data, list) and len(data) > 0 and not isinstance(data[0], dict):
            data = data[0]
    return data


def main(
    input_1=None,  # 模板路径
    input_2=None,  # 替换计划
    input_3=None,  # 输出路径（可选）
    body=None,
    **kwargs
) -> Dict:
    """
    主函数 - 替换模板文本
    
    Args:
        input_1: {模板路径}
        input_2: {替换计划} - JSON 格式的替换指令
        input_3: {输出路径} - 可选
        body: Agent body
    
    Returns:
        dict: 处理结果
    """
    try:
        # 解析模板路径
        template_path = flatten_input(input_1)
        if not isinstance(template_path, str):
            template_path = str(template_path) if template_path else ""
        
        # 解析替换计划
        substitution_plan = input_2
        if isinstance(substitution_plan, str):
            substitution_plan = json.loads(substitution_plan)
        if isinstance(substitution_plan, list) and len(substitution_plan) == 1:
            if isinstance(substitution_plan[0], list):
                substitution_plan = substitution_plan[0]
        
        # 解析输出路径
        output_path = flatten_input(input_3)
        if output_path and not isinstance(output_path, str):
            output_path = str(output_path)
        
        return substitute_template_text(
            template_path=template_path,
            substitution_plan=substitution_plan,
            output_path=output_path,
            auto_fit=True,  # 启用自动收缩文本以适应形状
            adjust_font_size=False,  # 可选：智能调整字体大小
            body=body
        )
        
    except Exception as e:
        import traceback
        return {
            "status": "error",
            "error": f"处理失败: {str(e)}",
            "traceback": traceback.format_exc()
        }


if __name__ == "__main__":
    import io
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')
    
    if len(sys.argv) >= 3:
        template = sys.argv[1]
        plan_json = sys.argv[2]
        output = sys.argv[3] if len(sys.argv) > 3 else None
        
        plan = json.loads(plan_json)
        result = main(template, plan, output)
        print(json.dumps(result, ensure_ascii=False, indent=2))
    else:
        print("""
替换模板文本工具

用法:
    python 替换模板文本.py <模板文件> <替换计划JSON> [输出路径]

参数:
    模板文件: PPTX 模板路径
    替换计划JSON: 替换指令的 JSON 字符串
    输出路径: 可选，默认生成 xxx_substituted.pptx

替换计划格式:
[
  {
    "page_index": 0,
    "replacements": [
      {"shape_index": 10, "new_text": "新标题"},
      {"shape_index": 11, "new_text": "新副标题"}
    ]
  },
  {
    "page_index": 1,
    "replacements": [
      {"shape_id": 5, "new_text": "第二页内容"}
    ]
  }
]

说明:
- shape_index: 形状在页面中的索引（从0开始）
- shape_id: 形状的唯一ID（可从提取结果中获取）
- 两者可以任选其一，优先使用 shape_index

示例:
    python 替换模板文本.py template.pptx '[{"page_index":0,"replacements":[{"shape_index":10,"new_text":"Hello"}]}]'
""")

